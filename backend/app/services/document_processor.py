import os
import google.generativeai as genai
from typing import Dict, Any, List
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
from ..config import settings

# Configure Gemini
genai.configure(api_key=settings.GOOGLE_API_KEY)

class DocumentProcessor:
    """Process various document types and extract text"""
    
    @staticmethod
    def process_pdf(file_path: str) -> Dict[str, Any]:
        """Extract text from PDF - with Vision Fallback"""
        text = ""
        metadata = {}
        
        try:
            # 1. Try Standard Extraction (PyPDF2)
            reader = PdfReader(file_path)
            metadata = {
                "pages": len(reader.pages),
                "file_type": "pdf"
            }
            
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
            
            text = text.strip()
            
            # 2. Vision Fallback (if text is empty/short)
            if len(text) < 50:
                print("⚠️ PDF appears to be an image (low text count). Switching to Gemini Vision OCR...")
                return DocumentProcessor._extract_with_gemini_vision(file_path)
            
            return {
                "text": text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"⚠️ Standard PDF extraction failed: {str(e)}")
            print("🔄 Attempting Gemini Vision Fallback...")
            return DocumentProcessor._extract_with_gemini_vision(file_path)

    @staticmethod
    def _extract_with_gemini_vision(file_path: str) -> Dict[str, Any]:
        """Use Gemini Vision to transcribe scanned documents"""
        try:
            print(f"📤 Uploading file to Gemini for Vision OCR: {os.path.basename(file_path)}")
            
            # Upload file
            uploaded_file = genai.upload_file(file_path)
            
            # Use Gemini Flash for speed (using the working model from logs)
            model = genai.GenerativeModel("models/gemini-flash-latest")
            
            prompt = """
            You are a high-precision OCR engine for startup documents.
            1. Transcribe ALL text in this document exactly as it appears.
            2. Describe any charts, graphs, or visual data in detail (e.g., "[Chart: Revenue Growth 2023-2025 showing 300% increase]").
            3. Do not summarize; provide the full content.
            """
            
            # Generate content
            response = model.generate_content([prompt, uploaded_file])
            transcription = response.text
            
            print(f"✅ Gemini Vision extracted {len(transcription)} characters")
            
            return {
                "text": transcription,
                "metadata": {
                    "file_type": "pdf_scanned",
                    "ocr_engine": "gemini_vision_flash"
                }
            }
            
        except Exception as e:
            raise Exception(f"Gemini Vision OCR failed: {str(e)}")

    @staticmethod
    def process_docx(file_path: str) -> Dict[str, Any]:
        """Extract text from DOCX"""
        try:
            doc = DocxDocument(file_path)
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text])
            
            metadata = {
                "paragraphs": len(doc.paragraphs),
                "file_type": "docx"
            }
            
            return {
                "text": text.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            raise Exception(f"DOCX processing failed: {str(e)}")

    @staticmethod
    def process_pptx(file_path: str) -> Dict[str, Any]:
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for i, slide in enumerate(prs.slides):
                text += f"\n--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            metadata = {
                "slides": len(prs.slides),
                "file_type": "pptx"
            }
            
            return {
                "text": text.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            raise Exception(f"PPTX processing failed: {str(e)}")

    @staticmethod
    def process_xlsx(file_path: str) -> Dict[str, Any]:
        """Extract text from XLSX"""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            
            metadata = {
                "sheets": len(wb.sheetnames),
                "file_type": "xlsx"
            }
            
            return {
                "text": text.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            raise Exception(f"XLSX processing failed: {str(e)}")

    @classmethod
    def process_file(cls, file_path: str) -> Dict[str, Any]:
        """Process file based on extension"""
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        processors = {
            ".pdf": cls.process_pdf,
            ".docx": cls.process_docx,
            ".pptx": cls.process_pptx,
            ".xlsx": cls.process_xlsx,
        }
        
        if ext not in processors:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return processors[ext](file_path)


# Singleton instance
document_processor = DocumentProcessor()